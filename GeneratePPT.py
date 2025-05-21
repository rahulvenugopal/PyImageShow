# -*- coding: utf-8 -*-
"""
Created on Wed May 21 14:16:50 2025
- Show random pics with captions
@author: Rahul Venugopal and ChatGPT
"""

# Load libraries
import os
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
from textwrap import wrap

# === CONFIGURATION ===
image_folder = r"C:\Users\Admin\Desktop\PyImageShow\Images"
output_pptx = r"C:\Users\Admin\Desktop\PyImageShow\Images\random_images_presentation.pptx"
image_extensions = ['.jpg', '.jpeg', '.png', '.bmp']

frame_width = Inches(8) # Adjust the box into which images fit in
frame_height = Inches(6)
frame_left = Inches(0.5)    # centered on 10" wide slide
frame_top = Inches(0.5)

# === COLLECT & SHUFFLE IMAGES ===
images = [f for f in os.listdir(image_folder)
          if os.path.splitext(f)[1].lower() in image_extensions]
random.shuffle(images)

# === CREATE PRESENTATION ===
prs = Presentation()
blank_layout = prs.slide_layouts[6]

for img_file in images:
    slide = prs.slides.add_slide(blank_layout)
    img_path = os.path.join(image_folder, img_file)

    # === RESIZE IMAGE TO FIT FRAME ===
    with Image.open(img_path) as img:
        img_width, img_height = img.size
        img_aspect = img_width / img_height
        frame_aspect = frame_width.inches / frame_height.inches

        if img_aspect > frame_aspect:
            # Fit to width
            disp_width = frame_width
            disp_height = Inches(frame_width.inches / img_aspect)
        else:
            # Fit to height
            disp_height = frame_height
            disp_width = Inches(frame_height.inches * img_aspect)

        # Center inside frame
        img_left = frame_left + (frame_width - disp_width) / 2
        img_top = frame_top + (frame_height - disp_height) / 2

        slide.shapes.add_picture(img_path, img_left, img_top,
                                 width=disp_width, height=disp_height)

    # === ADD FILENAME CAPTION ===
    caption_left = Inches(0.5) # Adjust the text location here
    caption_top = Inches(6.5)
    caption_width = Inches(8) # Controls text box length
    caption_height = Inches(1.0)
    textbox = slide.shapes.add_textbox(caption_left, caption_top,
                                       caption_width, caption_height)
    text_frame = textbox.text_frame
    text_frame.clear()

    filename = os.path.splitext(img_file)[0]
    wrapped = "\n".join(wrap(filename, width=60))
    p = text_frame.paragraphs[0]
    p.text = wrapped
    p.font.name = 'Calibri'
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.alignment = PP_ALIGN.JUSTIFY

# === SAVE ===
prs.save(output_pptx)
